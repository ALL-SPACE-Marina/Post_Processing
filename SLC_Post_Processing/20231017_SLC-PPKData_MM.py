from __future__ import print_function
import collections.abc as collections #DO NOT DELETE; required for importing the pptx module
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
import argparse
from datetime import date
from pptx.util import Inches, Pt
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt; plt.rcParams['font.size'] = 15
import os
import pip
plt.close('all')

# load
#fileName = r'G:\measurements\TLM_Meas\DVT_I-Type\Tx_TLM\QR00239\RHCP\B1\Cal_Investigation\SLC3\PP_Data\ppk_data'
fileName = r'C:\Terminal_Testing\Plots\ppk_data'
savePath=r'C:\Terminal_Testing\Plots'
plotsPath=r'C:\Terminal_Testing\Plots\Ass_Cal_Probe_Height_Comparison_Plots'
input_prs= r'C:\Terminal_Testing\Post_Processing\pptx_template.pptx' # location of the pptx template
pptxFname = 'Cal_Comp_Probe_Height.pptx' #text to be added to the filename of the pptx
pptxType= 'Calibration Comparison Probe Height' # text to be added to the title slide
beam='B1'

Th_deg_list=[0.0, 10.0, 20.0, 30.0, 40.0, 50.0, 55.0, 60.0, 65.0, 70.0] # measured theta pointing angles
Ph_deg=[0.0]#, 14.4, 28.8, 43.2, 57.6, 72.0, 86.4] #measured phi pointing angles
meas_3D='False' #defines whether the measurement is partial, full 3D or a single cut. If a single cut, value should be 'False', otherwisde 'True

TLM_Type='Tx'
Pol='RHCP'
QRcode='440-0329-00239'
PCBAver='175-0182/ Rev. 1.1'
meas_Type='TLM' #depending on what's been measured, this variable can have the following values: 'TLM' - measurement is performed only on a TLM level
                                                                                                # 'Single Lens' - measurements are performed only for the 3 lenses individually
                                                                                                # 'TLM and Single Lens' - measurements are performed for the TLM and each lens individually
                                                                                                # 'Lens 1', 'Lens 2' or 'Lens 3' - measurements are performed only for the specified individual lens

if TLM_Type=='Rx':
    cal_freq_list = [17.7, 18.2, 18.7, 19.2, 19.7, 20.2, 20.7, 21.2]
    meas_freq_list =[17.7, 18.2, 18.7, 19.2, 19.7, 20.2, 20.7, 21.2]
    f_min=17.2
    f_max=21.7
    y_min_scan=40
    y_max_scan=90
    y_min_scan_xpd = 0
    y_max_scan_xpd = 30
    slcNum=2
elif TLM_Type=='Tx':
    cal_freq_list=[27.5, 28.0, 28.5, 29.0, 29.5, 30.0, 30.5, 31.0]
    meas_freq_list=[27.5, 28.0, 28.5, 29.0, 29.5, 30.0, 30.5, 31.0]
    f_min = 27.0
    f_max = 31.5
    y_min_scan =20
    y_max_scan =80
    y_min_scan_xpd = 0
    y_max_scan_xpd = 30
    slcNum=3

if meas_Type=='TLM':
    lens_enab=['l1e_l2e_l3e']
    plot_title=['TLM']
elif meas_Type=='Single Lens':
    lens_enab=['l1e_l2d_l3d', 'l1d_l2e_l3d', 'l1d_l2d_l3d']
    plot_title = ['L1', 'L2', 'L3']
elif meas_Type=='TLM and Single Lens':
    lens_enab=['l1e_l2e_l3e', 'l1e_l2d_l3d', 'l1d_l2e_l3d', 'l1d_l2d_l3d']
    plot_title = ['TLM', 'L1', 'L2', 'L3']
elif meas_Type=='Lens 1':
    lens_enab='l1e_l2d_l3d'
    plot_title = ['L1']
elif meas_Type=='Lens 2':
    lens_enab='l1d_l2e_l3d'
    plot_title = ['L2']
elif meas_Type=='Lens 3':
    lens_enab='l1d_l2d_l3d'
    plot_title = ['L3']

# labelLog = ['QR00057_F-Type', 'QR00012_G-Type', 'QR00002_G-Type', 'QR00034_G-Type_1p5', 'QR00060_G-Type_1p5']
# fpathLog = [r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_G-Type_1p5/Tx_TLM/Comparisons/SLC3/TLM_Type_Comp_att0/Raw_Data/QR00057_F-Type',
#             r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_G-Type_1p5/Tx_TLM/Comparisons/SLC3/TLM_Type_Comp_att0/Raw_Data/QR00012_G-Type',
#             r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_G-Type_1p5/Tx_TLM/Comparisons/SLC3/TLM_Type_Comp_att0/Raw_Data/QR00002_G-Type',
#             r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_G-Type_1p5/Tx_TLM/Comparisons/SLC3/TLM_Type_Comp_att0/Raw_Data/QR00034_G-Type_1p5',
#             r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_G-Type_1p5/Tx_TLM/Comparisons/SLC3/TLM_Type_Comp_att0/Raw_Data/QR00060_G-Type_1p5']

# labels and identifying data from input file
labelLog = ['Ass_Cal_13mm', 'BB_Cal'] #'Cal_10mm', 'Cal_15mm', 'Cal_20mm']
fpathLog = [r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Tx_TLM/QR00239/RHCP/B1/Cal_Investigation/SLC3/Raw_Data/Ass_Cal',
            #r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Tx_TLM/QR00239/RHCP/B1/Cal_Investigation/SLC3/Raw_Data/Ass_Cal_10mm',
            #r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Tx_TLM/QR00239/RHCP/B1/Cal_Investigation/SLC3/Raw_Data/Ass_Cal_15mm',
            #r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Rx_TLM/Comparisons/SLC2/F_vs_G_vs_I/Raw_Data/QR00011_I-Type',
            r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Tx_TLM/QR00239/RHCP/B1/Cal_Investigation/SLC3/Raw_Data/BB_Cal']

# labelLog = ['Aluminum', 'Plastic'] #'QR00002_I-Type', 'QR00011_I-Type',
# fpathLog = [r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_F-Type/Rx_TLM/QR440-0111-00005/RHCP/B1/Nest_Test/SLC2/Raw_Data/Aluminum',
#             #r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Rx_TLM/Comparisons/SLC2/F_vs_G_vs_I/Raw_Data/QR00020_G-Type',
#             #r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Rx_TLM/Comparisons/SLC2/F_vs_G_vs_I/Raw_Data/QR00002_I-Type',
#             #r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_I-Type/Rx_TLM/Comparisons/SLC2/F_vs_G_vs_I/Raw_Data/QR00011_I-Type',
#             r'/mnt/nfs/data/groups/measurements/TLM_Meas/DVT_F-Type/Rx_TLM/QR440-0111-00005/RHCP/B1/Nest_Test/SLC2/Raw_Data/Plastic']


#defining markes and colors for the plots
markerList=['s', 'o', '^', 'D', 'X', 'p', '*', 'H']
colourMap = [['b','orange','g','r','purple','brown','magenta','grey'],
             ['c', 'peru', 'darkgreen', 'darksalmon', 'plum', 'chocolate', 'hotpink', 'k'],
             ['cornflowerblue', 'tan', 'greenyellow', 'darkred', 'blueviolet', 'rosybrown', 'mediumvioletred', 'darkslategrey'],
             ['deepskyblue', 'gold', 'olivedrab', 'crimson', 'indigo', 'sienna', 'orchid', 'silver']]

## Defs ##

def plot__paramVsScan(df, fpath, cal_freq, meas_freq, phi_deg, lens_enable):

    global dfPlot, x_scan_request, y_scan_request, xpd_scan_request, acu, phi_scan_mispoint, th_scan_mispoint

    # extracts information from the ppk file and calculates gain, XP, XPD and mispoint in theta and phi with respect to the requested pointing angle. The parameters are extracted for plots vs scan

    # reduce the dataframe size based on the available parameters for plotting vs frequency

    df = df[(df["pa_phi_deg"] == phi_deg)]
    df = df[(df["cal_freq_GHz"] == cal_freq)]
    df = df[(df["freq_GHz"] == meas_freq)]
    df = df[(df["fpath_parent"] == fpath)]
    df = df[(df["lens_enabled"] == lens_enable)]
    acu=np.array(df['acu_version'])[0]

    df_peak = df[(df["entry_type"] == 'gain_at_peak')]
    df_request = df[(df["entry_type"] == 'gain_at_requested_angle')]

    x_scan_request = np.array(df_request['theta_deg'])
    y_scan_request = np.array(df_request['Gain_dB'])
    xpd_scan_request=np.array(df_request['xpd_dB'])

    th_request = np.array(df_request['theta_deg'])
    th_peak = np.array(df_peak['theta_deg'])

    ph_peak = np.array(df_peak['phi_deg'])
    ph_request = np.array(df_request['phi_deg'])

    phi_scan_mispoint = ph_peak - ph_request
    th_scan_mispoint = th_peak - th_request


def plot__paramVsfreq(df, fpath,Th_deg,Ph_degree, lens_enable):
    global dfPlot, x_fq_request, y_fq_request, phi_meas_freq, xpd_freq_request, acu, phi_fq_mispoint, th_fq_mispoint, temp

    #extracts information from the ppk file and calculates gain, XP, XPD and mispoint in theta and phi with respect to the requested pointing angle. The parameters are extracted for plots vs frequency

    #reduce the dataframe size based on the available parameters for plotting vs frequency
    df = df[(df["pa_phi_deg"] == Ph_degree)]
    df = df[(df["pa_theta_deg"] == Th_deg)]
    df = df[(df["freq_GHz"] == df["cal_freq_GHz"])]
    df = df[(df["lens_enabled"] == lens_enable)]
    df = df[(df["fpath_parent"] == fpath)]
    acu = np.array(df['acu_version'])[0]
    temp = np.array(df['temp_C'])[0]

    df_peak = df[(df["entry_type"] == 'gain_at_peak')]
    df_request = df[(df["entry_type"] == 'gain_at_requested_angle')]

    th_request = np.array(df_request['theta_deg'])
    th_peak = np.array(df_peak['theta_deg'])

    ph_peak = np.array(df_peak['phi_deg'])
    ph_request = np.array(df_request['phi_deg'])

    phi_fq_mispoint = ph_peak - ph_request
    th_fq_mispoint = th_peak - th_request

    phi_meas_freq = np.array(df_request['pa_phi_deg'])
    x_fq_request = np.array(df_request['freq_GHz'])
    y_fq_request = np.array(df_request['Gain_dB'])
    xpd_freq_request = np.array(df_request['xpd_dB'])



# load the data from the xlsx input file
dirScript = os.getcwd()
os.chdir(dirScript)
df = pd.read_excel(fileName + ".xlsx")
columns = df.columns.tolist()

# for m in range(len(lens_enab)):
#
#     for k in range(len(cal_freq_list)):
#
#
#         for l in range(len(Ph_deg)):
#             # create figures
#             fig1 = plt.figure(figsize=([7, 6]))
#             ax1 = fig1.add_subplot(1, 1, 1)
#
#             fig2 = plt.figure(figsize=([7, 6]))
#             ax2 = fig2.add_subplot(1, 1, 1)
#
#             fig3 = plt.figure(figsize=([7, 6]))
#             ax3 = fig3.add_subplot(1, 1, 1)
#
#             if meas_3D == 'True':
#                 fig4 = plt.figure(figsize=([7, 6]))
#                 ax4 = fig4.add_subplot(1, 1, 1)
#
#             for i in range(len(fpathLog)):
#                 Ph_degree=Ph_deg[l]
#                 fpath = fpathLog[i]
#                 plot__paramVsScan(df, fpath, cal_freq_list[k], meas_freq_list[k], Ph_degree, lens_enab[m])
#
#
#                 # Plot Gain & XP
#                 ax1.plot(x_scan_request, y_scan_request, markerList[i], markerfacecolor= colourMap[0][i], markeredgecolor='k', markersize=10, label='Gain ' + labelLog[i])
#                 ax1.plot(x_scan_request, y_scan_request-xpd_scan_request, markerList[i], markerfacecolor= colourMap[2][i], markeredgecolor='k', markersize=10, label='XP ' + labelLog[i])
#                 ax1.set_xlabel('Theta, [deg]', fontsize=10);
#                 ax1.set_ylabel('Gain & XP, [dB]\n(at req. angle)', fontsize=10)
#                 yticks_param = np.linspace(0, 100, num=21)
#                 ax1.set_yticks(yticks_param)
#                 ax1.tick_params(axis='y', labelsize=10)
#                 ax1.tick_params(axis='x', labelsize=10)
#                 ax1.set_xlim([-10, 80]);
#                 ax1.set_ylim([y_min_scan, y_max_scan])
#                 ax1.grid('on')
#                 ax1.legend(loc='lower left', fontsize=10)
#                 ax1.set_title(
#                     str(plot_title[m]) + ' ' + Pol + ' Gain & XP'+'\nSW: ' + acu + '\nFreq = ' + str(meas_freq_list[k]) + ' GHz' + '\nTh=' + 'X' + ', Phi=' + str(Ph_degree),
#                     fontsize=15)
#                 fig1.tight_layout()
#                 fig1.savefig(
#                     savePath + '\\' + plot_title[m] + '_Gain_XP_vs_scan_Phi_' + str(Ph_degree) + 'Freq_' + str(meas_freq_list[k]) + 'GHz.png',
#                     dpi=400)
#
#
#                 #Plot Theta Mispoint
#                 ax2.plot(x_scan_request, th_scan_mispoint, markerList[i], markerfacecolor=colourMap[0][i],
#                          markeredgecolor='k', markersize=10, label='Gain ' + labelLog[i])
#                 ax2.set_xlabel('Theta, [deg]', fontsize=10);
#                 ax2.set_ylabel('Theta Pointing Error, [deg]', fontsize=10)
#                 yticks_param2 = np.linspace(-10, 10, num=21)
#                 ax2.set_yticks(yticks_param2)
#                 ax2.tick_params(axis='y', labelsize=10)
#                 ax2.tick_params(axis='x', labelsize=10)
#                 ax2.set_xlim([-10, 80]);
#                 ax2.set_ylim([-10, 10]);
#                 ax2.grid('on')
#                 ax2.legend(loc='lower left', fontsize=10)
#                 ax2.set_title(str(
#                     plot_title[m]) + ' ' + Pol + ' Theta Mispoint \nSW: ' + acu + '\nFreq = ' + str(meas_freq_list[k]) + ' GHz' + '\nTh=' + 'X' + ', Phi=' + str(Ph_degree), fontsize=15)
#                 fig2.tight_layout()
#                 fig2.savefig(savePath + '\\' + str(plot_title[m]) + '_Theta_Mispoint_vs_Scan_Ph_' + str(Ph_degree) + 'Freq_' + str(meas_freq_list[k]) + 'GHz.png', dpi=400)
#
#
#                 # Plot XPD
#                 ax3.plot(x_scan_request, xpd_scan_request, markerList[i], markerfacecolor=colourMap[0][i],
#                          markeredgecolor='k', markersize=10, label='XPD ' + labelLog[i])
#                 ax3.set_xlabel('Theta, [deg]', fontsize=10);
#                 ax3.set_ylabel('XPD, [dB]\n(at req. angle)', fontsize=10)
#                 yticks_param3 = np.linspace(0, 100, num=21)
#                 ax3.set_yticks(yticks_param3)
#                 ax3.tick_params(axis='y', labelsize=10)
#                 ax3.tick_params(axis='x', labelsize=10)
#                 ax3.set_xlim([-10, 80]);
#                 ax3.set_ylim([0, 30]);
#                 ax3.grid('on')
#                 ax3.legend(loc='lower left', fontsize=10)
#                 ax3.set_title(
#                     str(plot_title[m]) + ' ' + Pol + ' XPD' + '\nSW: ' + acu + '\nFreq = ' + str(
#                         meas_freq_list[k]) + ' GHz' + '\nTh=' + 'X' + ', Phi=' + str(Ph_degree),
#                     fontsize=15)
#                 fig1.tight_layout()
#                 fig1.savefig(
#                     savePath + '\\' + plot_title[m] + '_XPD_vs_scan_Phi_' + str(Ph_degree) + 'Freq_' + str(
#                         meas_freq_list[k]) + 'GHz.png',
#                     dpi=400)
#
#
#                 # Plot Phi Mispoint
#                 if meas_3D=='True':
#                     ax4.plot(x_scan_request, phi_scan_mispoint, markerList[i], markerfacecolor=colourMap[0][i],
#                              markeredgecolor='k', markersize=10, label='Gain ' + labelLog[i])
#                     ax4.set_xlabel('Theta, [deg]', fontsize=10);
#                     ax4.set_ylabel('Phi Pointing Error, [deg]', fontsize=10)
#                     yticks_param4 = np.linspace(-10, 10, num=21)
#                     ax4.set_yticks(yticks_param4)
#                     ax4.tick_params(axis='y', labelsize=10)
#                     ax4.tick_params(axis='x', labelsize=10)
#                     ax4.set_xlim([-10, 80]);
#                     ax4.set_ylim([-10, 10]);
#                     ax4.grid('on')
#                     ax4.legend(loc='lower left', fontsize=10)
#                     ax4.set_title(str(
#                         plot_title[m]) + ' ' + Pol + ' Phi Mispoint \nSW: ' + acu + '\nFreq = ' + str(
#                         meas_freq_list[k]) + ' GHz' + '\nTh=' + 'X' + ', Phi=' + str(Ph_degree), fontsize=15)
#                     fig4.tight_layout()
#                     fig4.savefig(savePath + '\\' + str(plot_title[m]) + '_Phi_Mispoint_vs_Scan_Ph_' + str(
#                         Ph_degree) + 'Freq_' + str(meas_freq_list[k]) + 'GHz.png', dpi=400)
#
#             plt.close('all')
#
#
#     for p in range(len(Th_deg_list)):
#
#         for l in range(len(Ph_deg)):
#
#             #create figures
#             fig1=plt.figure(figsize=([7, 6]))
#             ax1=fig1.add_subplot(1,1,1)
#
#             fig2 = plt.figure(figsize=([7, 6]))
#             ax2 = fig2.add_subplot(1, 1, 1)
#
#             fig3 = plt.figure(figsize=([7, 6]))
#             ax3 = fig3.add_subplot(1, 1, 1)
#
#             if meas_3D=='True':
#                 fig4 = plt.figure(figsize=([7, 6]))
#                 ax4 = fig4.add_subplot(1, 1, 1)
#
#             for i in range(len(fpathLog)):
#                 Ph_degree=Ph_deg[l]
#                 fpath = fpathLog[i]
#                 plot__paramVsfreq(df, fpath, Th_deg_list[p], Ph_degree, lens_enab[m])
#
#
#                 #plot Gain and XP
#                 ax1.plot(x_fq_request, y_fq_request, markerList[i], markerfacecolor= colourMap[0][i], markeredgecolor='k', markersize=10, label='Gain ' + labelLog[i])
#                 ax1.plot(x_fq_request, y_fq_request - xpd_freq_request, markerList[i], markerfacecolor=colourMap[2][i], markeredgecolor='k', markersize=10, label='XP ' + labelLog[i])
#                 ax1.set_xlabel('Frequency, [GHz]', fontsize=10);
#                 ax1.set_ylabel('Gain & XP, [dB]\n(at req. angle)', fontsize=10)
#                 yticks_param=np.linspace(0,100,num=21)
#                 ax1.set_yticks(yticks_param)
#                 ax1.tick_params(axis='y', labelsize=10)
#                 ax1.tick_params(axis='x', labelsize=10)
#                 ax1.set_xlim([f_min,f_max]);
#                 ax1.set_ylim([y_min_scan,y_max_scan]);# plt.ylim([65, 75])
#                 ax1.grid('on')
#                 ax1.legend(loc='lower left', fontsize=10)
#                 ax1.set_title(str(plot_title[m]) + ' ' + Pol + 'Gain & XP \nSW: ' + acu + '\nFreq = ' + ' X' + '\nTh=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree), fontsize=15)
#                 fig1.tight_layout()
#                 fig1.savefig(savePath + '\\' + str(plot_title[m]) + '_Gain_XP_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png', dpi=400)
#
#
#                 #Plot Theta Mispoint
#                 ax2.plot(x_fq_request, th_fq_mispoint, markerList[i], markerfacecolor= colourMap[0][i], markeredgecolor='k', markersize=10, label='Gain ' + labelLog[i])
#                 ax2.set_xlabel('Frequency, [GHz]', fontsize=10);
#                 ax2.set_ylabel('Pointing Error, [deg]', fontsize=10)
#                 yticks_param2 = np.linspace(-10, 10, num=21)
#                 ax2.set_yticks(yticks_param2)
#                 ax2.tick_params(axis='y', labelsize=10)
#                 ax2.tick_params(axis='x', labelsize=10)
#                 ax2.set_xlim([f_min,f_max]);
#                 ax2.set_ylim([-10,10]);
#                 ax2.grid('on')
#                 ax2.legend(loc='lower left', fontsize=10)
#                 ax2.set_title(str(plot_title[m]) + ' ' + Pol + ' Theta Mispoint \nSW: ' + acu + '\nFreq = ' + ' X' + '\nTh=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree), fontsize=15)
#                 fig2.tight_layout()
#                 fig2.savefig(savePath + '\\' + str(plot_title[m]) + '_Theta_Mispoint_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png', dpi=400)
#
#
#                 # plot XPD
#                 ax3.plot(x_fq_request, xpd_freq_request, markerList[i], markerfacecolor=colourMap[0][i], markeredgecolor='k',
#                          markersize=10, label='XPD ' + labelLog[i])
#                 ax3.set_xlabel('Frequency, [GHz]', fontsize=10);
#                 ax3.set_ylabel('XPD, [dB]\n(at req. angle)', fontsize=10)
#                 yticks_param3 = np.linspace(0, 30, num=7)
#                 ax3.set_yticks(yticks_param3)
#                 ax3.tick_params(axis='y', labelsize=10)
#                 ax3.tick_params(axis='x', labelsize=10)
#                 ax3.set_xlim([f_min, f_max]);
#                 ax3.set_ylim([0, 30]);  # plt.ylim([65, 75])
#                 ax3.grid('on')
#                 ax3.legend(loc='lower left', fontsize=10)
#                 ax3.set_title(
#                     str(plot_title[m]) + ' ' + Pol + 'XPD \nSW: ' + acu + '\nFreq = ' + ' X' + '\nTh=' + str(
#                         Th_deg_list[p]) + ', Phi=' + str(Ph_degree), fontsize=15)
#                 fig3.tight_layout()
#                 fig3.savefig(savePath + '\\' + str(plot_title[m]) + '_XPD_vs_Frequency_Th_' + str(
#                     Th_deg_list[p]) + '_Phi_' + str(Ph_degree) + '.png', dpi=400)
#
#
#                 #plot Phi mispointing if more than 1 phi cut has been measured (partial 3D measurement)
#                 if meas_3D=='True':
#                     ax4.plot(x_fq_request, ph_fq_mispoint, markerList[i], markerfacecolor=colourMap[0][i],
#                              markeredgecolor='k', markersize=10, label='Gain ' + labelLog[i])
#                     ax4.set_xlabel('Frequency, [GHz]', fontsize=10);
#                     ax4.set_ylabel('Phi Pointing Error, [deg]', fontsize=10)
#                     yticks_param4 = np.linspace(-10, 10, num=21)
#                     ax4.set_yticks(yticks_param2)
#                     ax4.tick_params(axis='y', labelsize=10)
#                     ax4.tick_params(axis='x', labelsize=10)
#                     ax4.set_xlim([f_min, f_max]);
#                     ax4.set_ylim([-10, 10]);
#                     ax4.grid('on')
#                     ax4.legend(loc='lower left', fontsize=10)
#                     ax4.set_title(str(plot_title[
#                                           m]) + ' ' + Pol + ' Phi Mispoint \nSW: ' + acu + '\nFreq = ' + ' X' + '\nTh=' + str(
#                         Th_deg_list[p]) + ', Phi=' + str(Ph_degree), fontsize=15)
#                     fig4.tight_layout()
#                     fig4.savefig(savePath + '\\' + str(plot_title[m]) + '_Phi_Mispoint_vs_Frequency_Th_' + str(
#                         Th_deg_list[p]) + '_Phi_' + str(Ph_degree) + '.png', dpi=400)
#
#             plt.close('all')
#

prs = Presentation(input_prs)

#define layout of slides
title_slide_layout = prs.slide_layouts[1]
overview_slide_layout = prs.slide_layouts[8]
measSet_slide_layout = prs.slide_layouts[29]
break_slide_layout = prs.slide_layouts[24]
data_slide_layout = prs.slide_layouts[29]
break2_slide_layout = prs.slide_layouts[1]
concl_slide_layout = prs.slide_layouts[9]
app_slide_layout= prs.slide_layouts[9]


#create title slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = pptxType
subtitle.text = TLM_Type+' '+ meas_Type+ ' QR Code: '+QRcode+ '\nDate: '+str(date.today())

#crete overview slide
slide = prs.slides.add_slide(overview_slide_layout)
title = slide.shapes.title
title.text = 'Overview'

placeholder=slide.placeholders[13]
p=placeholder.text_frame.add_paragraph()
run=p.add_run()
p.level=1
run.text= 'Measurement campaign details'

p=placeholder.text_frame.add_paragraph()
run1=p.add_run()
p.level=2
run1.text= 'Measurements were performed in SLC' + str(slcNum)+\
          '\n '\
           '\n'

p=placeholder.text_frame.add_paragraph()
run2=p.add_run()
p.level=1
run2.text= '\n' \
           'Main observations'

p=placeholder.text_frame.add_paragraph()
run1=p.add_run()
p.level=2
run1.text= '' \
          '\n '\
           '\n'

#create measurement config slide
slide = prs.slides.add_slide(measSet_slide_layout)
title = slide.shapes.title
title.text = 'Measurement Settings'

x, y, cx, cy = Inches(0.5), Inches(2), Inches(10), Inches(1)
shape = slide.shapes.add_table(2, 5, x, y, cx, cy)
style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'

tbl =  shape._element.graphic.graphicData.tbl
tbl[0][-1].text = style_id

table=shape.table   #design#(style_id)


cell11=table.cell(0,0)
cell11.text= 'TLM Type'

cell12=table.cell(0,1)
cell12.text= 'TLM PCBA Version'

cell13=table.cell(0,2)
cell13.text= 'TLM Ass QR Code '

cell14=table.cell(0,3)
cell14.text= 'TSW Version'

cell15=table.cell(0,4)
cell15.text= 'Measurement Temperature'

acu='1.20.26'
temp='45'
table.cell(1,0).text= str(TLM_Type)
table.cell(1,1).text= str(PCBAver)
table.cell(1,2).text= str(QRcode)
table.cell(1,3).text= str(acu)
table.cell(1,4).text= str(temp)

# placeholder=slide.placeholders[0]
# shape=placeholder.table()   #insert_table(rows=3, cols=4)
# table=shape.table



#add plots slide
for m in range(len(lens_enab)):
    # create break slide
    slide = prs.slides.add_slide(break_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TLM_Type + ' ' + meas_Type + ' Gain and XP performance vs Scan '
    subtitle.text = str(Pol) + ', ' + str(beam)

    # add plots slide
    for m in range(len(lens_enab)):
        for k in range(len(cal_freq_list)):
            for l in range(len(Ph_deg)):
                Ph_degree = Ph_deg[l]

                slide = prs.slides.add_slide(data_slide_layout)
                title = slide.shapes.title
                title.text = 'Freq = ' + str(meas_freq_list[k]) + ' GHz, ' + 'Th=' + 'X' + ', Phi=' + str(Ph_degree)
                input_plt = plotsPath + '\\' + plot_title[m] + '_Gain_XP_vs_scan_Phi_' + str(Ph_degree) + 'Freq_' + str(
                    meas_freq_list[k]) + 'GHz.png'
                plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))


    slide = prs.slides.add_slide(break_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TLM_Type + ' ' + plot_title[m] + ' XPD vs Scan '
    subtitle.text = str(Pol) + ', ' + str(beam)

    for k in range(len(cal_freq_list)):
        for l in range(len(Ph_deg)):
            Ph_degree=Ph_deg[l]

            slide = prs.slides.add_slide(data_slide_layout)
            title = slide.shapes.title
            title.text = 'Freq = ' + str(meas_freq_list[k]) + ' GHz, ' + 'Th=' + 'X' + ', Phi=' + str(Ph_degree)
            input_plt= plotsPath + '\\' + plot_title[m] + '_XPD_vs_scan_Phi_' + str(Ph_degree) + 'Freq_' + str(meas_freq_list[k]) + 'GHz.png'
            plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))

    slide = prs.slides.add_slide(break_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TLM_Type + ' ' + plot_title[m] + ' Theta Mispoint vs Scan '
    subtitle.text = str(Pol) + ', ' + str(beam)

    for k in range(len(cal_freq_list)):
        for l in range(len(Ph_deg)):
            Ph_degree = Ph_deg[l]

            slide = prs.slides.add_slide(data_slide_layout)
            title = slide.shapes.title
            title.text = 'Freq = ' + str(meas_freq_list[k]) + ' GHz, ' + 'Th=' + 'X' + ', Phi=' + str(Ph_degree)
            input_plt = plotsPath + '\\' + str(plot_title[m]) + '_Theta_Mispoint_vs_Scan_Ph_' + str(Ph_degree) + 'Freq_' + str(meas_freq_list[k]) + 'GHz.png'
            plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))

    if meas_3D== 'True':
        slide = prs.slides.add_slide(break_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = TLM_Type + ' ' + plot_title[m] + ' Phi Mispoint vs Scan '
        subtitle.text = str(Pol) + ', ' + str(beam)

        for k in range(len(cal_freq_list)):
            for l in range(len(Ph_deg)):
                Ph_degree = Ph_deg[l]

                slide = prs.slides.add_slide(data_slide_layout)
                title = slide.shapes.title
                title.text = 'Freq = ' + str(meas_freq_list[k]) + ' GHz, ' + 'Th=' + 'X' + ', Phi=' + str(Ph_degree)
                input_plt = plotsPath + '\\' + str(plot_title[m]) + '_Phi_Mispoint_vs_Scan_Ph_' + str(
                    Ph_degree) + 'Freq_' + str(meas_freq_list[k]) + 'GHz.png'
                plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))



for m in range(len(lens_enab)):
    # create break slide
    slide = prs.slides.add_slide(break_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TLM_Type + ' ' + meas_Type + ' Gain and XP performance vs Frequency '
    subtitle.text = str(Pol) + ', ' + str(beam)

    #                 ax1.set_title(str(plot_title[m]) + ' ' + Pol + 'Gain & XP \nSW: ' + acu + '\nFreq = ' + ' X' + '\nTh=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree), fontsize=15)
    #                 fig1.tight_layout()
    #                 fig1.savefig(savePath + '\\' + str(plot_title[m]) + '_Gain_XP_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png', dpi=400)
    #

    for p in range(len(Th_deg_list)):
        for l in range(len(Ph_deg)):
            Ph_degree = Ph_deg[l]

            slide = prs.slides.add_slide(data_slide_layout)
            title = slide.shapes.title
            title.text = 'Freq = ' + ' X' + ', Th=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree)
            input_plt = plotsPath + '\\' + plot_title[m] + '_Gain_XP_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png'
            plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))


    slide = prs.slides.add_slide(break_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TLM_Type + ' ' + plot_title[m] + ' XPD vs Frequency '
    subtitle.text = str(Pol) + ', ' + str(beam)

    for p in range(len(Th_deg_list)):
        for l in range(len(Ph_deg)):
            Ph_degree=Ph_deg[l]

            slide = prs.slides.add_slide(data_slide_layout)
            title = slide.shapes.title
            title.text = 'Freq = ' + ' X' + ', Th=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree)
            input_plt= plotsPath + '\\' + plot_title[m] + '_XPD_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png'
            plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))

    slide = prs.slides.add_slide(break_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TLM_Type + ' ' + plot_title[m] + ' Theta Mispoint vs Frequency '
    subtitle.text = str(Pol) + ', ' + str(beam)

    for p in range(len(Th_deg_list)):
        for l in range(len(Ph_deg)):
            Ph_degree = Ph_deg[l]

            slide = prs.slides.add_slide(data_slide_layout)
            title = slide.shapes.title
            title.text = 'Freq = ' + ' X' + ', Th=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree)
            input_plt = plotsPath + '\\' + str(plot_title[m]) + '_Theta_Mispoint_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png'
            plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))

    if meas_3D== 'True':
        slide = prs.slides.add_slide(break_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = TLM_Type + ' ' + plot_title[m] + ' Phi Mispoint vs Frequency '
        subtitle.text = str(Pol) + ', ' + str(beam)

        for p in range(len(Th_deg_list)):
            for l in range(len(Ph_deg)):
                Ph_degree = Ph_deg[l]

                slide = prs.slides.add_slide(data_slide_layout)
                title = slide.shapes.title
                title.text = 'Freq = ' + ' X' + ', Th=' + str(Th_deg_list[p]) + ', Phi=' + str(Ph_degree)
                input_plt = plotsPath + '\\' + str(plot_title[m]) + '_Phi_Mispoint_vs_Frequency_Th_'+ str(Th_deg_list[p]) + '_Phi_' + str(Ph_degree)+'.png'
                plot = slide.shapes.add_picture(input_plt, Inches(1), Inches(1))


slide = prs.slides.add_slide(break2_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Conclusions'
subtitle.text = TLM_Type+' '+ meas_Type+ ' QR Code: '+QRcode


#crete conclusions slide
slide = prs.slides.add_slide(concl_slide_layout)
title = slide.shapes.title
title.text = 'Conclusions'

placeholder=slide.placeholders[13]
p=placeholder.text_frame.add_paragraph()
run=p.add_run()
p.level=1
run.text= 'Information on the test campaign'

p=placeholder.text_frame.add_paragraph()
run1=p.add_run()
p.level=2
run1.text= '\n '\
           '\n'

p=placeholder.text_frame.add_paragraph()
run2=p.add_run()
p.level=1
run2.text= '\n' \
           'Key elements of the measurement campaign'

p=placeholder.text_frame.add_paragraph()
run1=p.add_run()
p.level=2
run1.text= '' \
          '\n '\
           '\n'

p=placeholder.text_frame.add_paragraph()
run2=p.add_run()
p.level=1
run2.text= '\n' \
           'Key measurement results'

p=placeholder.text_frame.add_paragraph()
run1=p.add_run()
p.level=2
run1.text= '' \
          '\n '\
           '\n'






prs.save(savePath+'\\'+ pptxFname)

#parse_args('test.ppt, test-output-markup.ppt)
#analyze_ppt('test.pptx', 'test-output-markup.pptx')

#create_ppt(savePath+'\\'+'test1.pptx', savePath+'\\'+ 'test-output.pptx')